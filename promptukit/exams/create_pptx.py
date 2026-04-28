"""Generate a PowerPoint (.pptx) deck from a promptukit question bank.

Mirrors the CLI shape of ``create_exam.py`` but emits one slide per question
(plus an optional answer slide) using ``python-pptx``. A flag controls whether
answers appear on a separate following slide, inline at the bottom of the
question slide, or are omitted entirely.

Run: ``python -m promptukit.exams.create_pptx -q bank.json -o deck.pptx --answers after``
"""

from __future__ import annotations

import argparse
import json
import re
import sys
import unicodedata
from pathlib import Path
from typing import Any, Iterable, List, Optional

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Emu, Inches, Pt

from promptukit.exams.create_exam import (
    DEFAULT_METADATA,
    effective_metadata,
    load_metadata_from_json,
    questions_artifact,
    save_json_artifact,
)
from promptukit.questions.question_models import (
    Calculation,
    FillInTheBlank,
    Matching,
    MultipleChoice,
    Question,
    ShortAnswer,
    TrueFalse,
)


ANSWER_MODES = ("none", "after", "inline")

_TAG_RE = re.compile(r"<[^>]+>")
_ENTITY_MAP = {
    "&mdash;": "—",
    "&ndash;": "–",
    "&minus;": "−",
    "&nbsp;": " ",
    "&deg;": "°",
    "&Delta;": "Δ",
    "&phi;": "φ",
    "&pi;": "π",
    "&times;": "×",
    "&plusmn;": "±",
    "&amp;": "&",
    "&lt;": "<",
    "&gt;": ">",
    "&quot;": '"',
}
_LEADING_NUMBER_RE = re.compile(r"^\s*\d+\.\s*")
_LETTER_PREFIX_RE = re.compile(r"^[A-Za-z][\)\.\s]+")

MUTED = RGBColor(0x80, 0x80, 0x80)
DARK = RGBColor(0x1A, 0x1A, 0x1A)


def pptx_safe_text(value: Any) -> str:
    """Plain-text projection of a question string for python-pptx runs.

    Strips ReportLab inline tags (``<b>``, ``<sub>``, ...) since python-pptx
    formats text via run attributes, not markup. Resolves common HTML entities
    used in existing banks. Normalises to NFC.
    """
    s = str(value if value is not None else "")
    s = unicodedata.normalize("NFC", s)
    for ent, ch in _ENTITY_MAP.items():
        s = s.replace(ent, ch)
    s = _TAG_RE.sub("", s)
    return s


def _flatten_questions(sections_or_questions: Any) -> List[dict]:
    """Return a flat list of question dicts regardless of input shape."""
    if isinstance(sections_or_questions, dict):
        if "sections" in sections_or_questions:
            sections_or_questions = sections_or_questions["sections"]
        elif "questions" in sections_or_questions:
            return [q for q in sections_or_questions["questions"] if isinstance(q, dict)]
    if isinstance(sections_or_questions, list):
        if (
            sections_or_questions
            and isinstance(sections_or_questions[0], dict)
            and "questions" in sections_or_questions[0]
        ):
            flat: List[dict] = []
            for sec in sections_or_questions:
                for q in sec.get("questions") or []:
                    if isinstance(q, dict):
                        flat.append(q)
            return flat
        return [q for q in sections_or_questions if isinstance(q, dict)]
    return []


def _add_textbox(
    slide,
    left,
    top,
    width,
    height,
    text: str,
    *,
    size: int = 18,
    bold: bool = False,
    italic: bool = False,
    align=PP_ALIGN.LEFT,
    color: Optional[RGBColor] = None,
):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = pptx_safe_text(text)
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    if color is not None:
        run.font.color.rgb = color
    return box


def _add_bullets(slide, left, top, width, height, lines: Iterable[str], *, size: int = 20):
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = PP_ALIGN.LEFT
        run = p.add_run()
        run.text = pptx_safe_text(line)
        run.font.size = Pt(size)
    return box


def _slide_geometry(prs: Presentation):
    sw = prs.slide_width
    sh = prs.slide_height
    margin = Inches(0.5)
    inner_w = sw - 2 * margin
    return {
        "slide_w": sw,
        "slide_h": sh,
        "margin": margin,
        "inner_w": inner_w,
        "title_top": Inches(0.4),
        "title_h": Inches(1.0),
        "body_top": Inches(1.6),
        "body_h": Inches(5.0),
        "footer_top": sh - Inches(0.5),
        "footer_h": Inches(0.35),
        "number_w": Inches(1.5),
    }


def _add_footer(slide, geom, meta):
    parts = [meta.get("institution") or "", meta.get("instructor") or ""]
    text = " — ".join(p for p in parts if p)
    if not text:
        return
    _add_textbox(
        slide,
        geom["margin"],
        geom["footer_top"],
        geom["inner_w"],
        geom["footer_h"],
        text,
        size=10,
        color=MUTED,
        align=PP_ALIGN.CENTER,
    )


def _add_question_number(slide, geom, number: int, total: int):
    box = slide.shapes.add_textbox(
        geom["slide_w"] - geom["margin"] - geom["number_w"],
        geom["title_top"],
        geom["number_w"],
        Inches(0.4),
    )
    tf = box.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f"{number} / {total}"
    run.font.size = Pt(11)
    run.font.color.rgb = MUTED


def _add_title(slide, geom, text: str):
    _add_textbox(
        slide,
        geom["margin"],
        geom["title_top"],
        geom["inner_w"] - geom["number_w"],
        geom["title_h"],
        text,
        size=24,
        bold=True,
        color=DARK,
    )


def _strip_question_number(text: str) -> str:
    return _LEADING_NUMBER_RE.sub("", text or "", count=1).strip()


def _letter(idx: int) -> str:
    return chr(ord("A") + idx)


def _format_choice(idx: int, raw: str) -> str:
    raw = pptx_safe_text(str(raw)).strip()
    if _LETTER_PREFIX_RE.match(raw):
        return raw
    return f"{_letter(idx)}) {raw}"


# --- Per-type renderers --------------------------------------------------

def _render_multiple_choice(slide, geom, q: MultipleChoice):
    _add_title(slide, geom, _strip_question_number(q.text))
    lines = [_format_choice(i, c) for i, c in enumerate(q.choices)]
    _add_bullets(slide, geom["margin"], geom["body_top"], geom["inner_w"], geom["body_h"], lines)


def _render_truefalse(slide, geom, q: TrueFalse):
    _add_title(slide, geom, _strip_question_number(q.text))
    _add_textbox(
        slide,
        geom["margin"],
        geom["body_top"] + Inches(1.0),
        geom["inner_w"],
        Inches(2.0),
        "True        •        False",
        size=32,
        bold=True,
        align=PP_ALIGN.CENTER,
    )


def _render_short_answer(slide, geom, q: ShortAnswer):
    _add_title(slide, geom, _strip_question_number(q.text))


def _render_fill_in_blank(slide, geom, q: FillInTheBlank):
    _add_title(slide, geom, _strip_question_number(q.text))


def _render_matching_table(slide, geom, q: Matching, *, filled: bool):
    pairs = q.pairs or []
    rows = max(len(pairs), 1)
    table_shape = slide.shapes.add_table(
        rows,
        2,
        geom["margin"],
        geom["body_top"],
        geom["inner_w"],
        min(geom["body_h"], Emu(int(0.6 * 914400) * rows)),
    )
    table = table_shape.table
    for i, pair in enumerate(pairs):
        left = pair[0] if len(pair) > 0 else ""
        right = pair[1] if len(pair) > 1 else ""
        for j, val in enumerate((left, right)):
            cell = table.cell(i, j)
            cell.text = ""
            tf = cell.text_frame
            p = tf.paragraphs[0]
            run = p.add_run()
            run.text = pptx_safe_text(val)
            run.font.size = Pt(18)
            if filled and j == 1:
                run.font.bold = True


def _render_matching(slide, geom, q: Matching):
    _add_title(slide, geom, _strip_question_number(q.text))
    _render_matching_table(slide, geom, q, filled=False)


def _render_calculation(slide, geom, q: Calculation):
    _add_title(slide, geom, _strip_question_number(q.text))
    if q.unit:
        _add_textbox(
            slide,
            geom["margin"],
            geom["body_top"] + Inches(0.4),
            geom["inner_w"],
            Inches(0.5),
            f"Units: {q.unit}",
            size=14,
            italic=True,
            color=MUTED,
        )


_RENDERERS = {
    MultipleChoice: _render_multiple_choice,
    TrueFalse: _render_truefalse,
    ShortAnswer: _render_short_answer,
    FillInTheBlank: _render_fill_in_blank,
    Matching: _render_matching,
    Calculation: _render_calculation,
}


# --- Answer rendering ----------------------------------------------------

def _answer_text(q: Question) -> Optional[str]:
    if isinstance(q, MultipleChoice):
        if q.answer_index is not None and 0 <= q.answer_index < len(q.choices):
            choice = pptx_safe_text(q.choices[q.answer_index])
            choice = _LETTER_PREFIX_RE.sub("", choice)
            return f"{_letter(q.answer_index)}) {choice}"
        if q.answer_text is not None:
            return pptx_safe_text(q.answer_text)
        return None
    if isinstance(q, TrueFalse):
        if q.answer is None:
            return None
        return "True" if q.answer else "False"
    if isinstance(q, ShortAnswer):
        return q.answer or None
    if isinstance(q, FillInTheBlank):
        if not q.answers:
            return None
        return ", ".join(q.answers)
    if isinstance(q, Calculation):
        if q.answer is None:
            return None
        body = f"{q.answer}"
        if q.unit:
            body = f"{body} {q.unit}"
        if q.tolerance:
            body = f"{body} ± {q.tolerance}"
        return body
    return None


def _render_answer_body(slide, geom, q: Question):
    if isinstance(q, FillInTheBlank):
        prompt = _strip_question_number(q.text)
        filled = prompt
        for ans in q.answers:
            filled = re.sub(r"_{2,}", pptx_safe_text(ans), filled, count=1)
        _add_textbox(
            slide,
            geom["margin"],
            geom["body_top"],
            geom["inner_w"],
            geom["body_h"],
            filled,
            size=22,
            bold=True,
        )
        return
    if isinstance(q, Matching):
        _render_matching_table(slide, geom, q, filled=True)
        return
    text = _answer_text(q) or "(no answer provided)"
    _add_textbox(
        slide,
        geom["margin"],
        geom["body_top"],
        geom["inner_w"],
        geom["body_h"],
        text,
        size=28,
        bold=True,
    )


def _add_answer_slide(prs, geom, meta, q: Question):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_title(slide, geom, "Answer")
    _render_answer_body(slide, geom, q)
    _add_footer(slide, geom, meta)
    return slide


def _add_inline_answer(slide, geom, q: Question):
    text = _answer_text(q)
    if isinstance(q, FillInTheBlank) and q.answers:
        text = ", ".join(q.answers)
    if isinstance(q, Matching) and q.pairs:
        text = "; ".join(f"{p[0]} → {p[1]}" for p in q.pairs if len(p) == 2)
    if not text:
        return
    _add_textbox(
        slide,
        geom["margin"],
        geom["footer_top"] - Inches(0.6),
        geom["inner_w"],
        Inches(0.5),
        f"Answer: {text}",
        size=14,
        italic=True,
        color=MUTED,
    )


# --- Cover & question slides --------------------------------------------

def _add_cover_slide(prs, geom, meta):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    title = pptx_safe_text(meta.get("title") or "Question Bank")
    _add_textbox(
        slide,
        geom["margin"],
        Inches(2.0),
        geom["inner_w"],
        Inches(1.5),
        title,
        size=44,
        bold=True,
        align=PP_ALIGN.CENTER,
        color=DARK,
    )
    subtitle_parts = [meta.get("institution") or "", meta.get("instructor") or ""]
    subtitle = " — ".join(p for p in subtitle_parts if p)
    if subtitle:
        _add_textbox(
            slide,
            geom["margin"],
            Inches(3.6),
            geom["inner_w"],
            Inches(0.6),
            subtitle,
            size=20,
            align=PP_ALIGN.CENTER,
            color=MUTED,
        )
    exam_type = meta.get("exam_type")
    if exam_type:
        _add_textbox(
            slide,
            geom["margin"],
            Inches(4.4),
            geom["inner_w"],
            Inches(0.5),
            exam_type,
            size=16,
            italic=True,
            align=PP_ALIGN.CENTER,
            color=MUTED,
        )
    return slide


def _render_question_slide(prs, geom, meta, q: Question, number: int, total: int, answers: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    renderer = _RENDERERS.get(type(q))
    if renderer is None:
        _add_title(slide, geom, _strip_question_number(q.text))
    else:
        renderer(slide, geom, q)
    _add_question_number(slide, geom, number, total)
    _add_footer(slide, geom, meta)
    if answers == "inline":
        _add_inline_answer(slide, geom, q)
    return slide


# --- JSON loading --------------------------------------------------------

def _read_question_bank(path: str) -> Any:
    """Load a JSON bank preserving full per-question fields.

    Supports the same top-level layouts as ``create_exam.load_questions_from_json``
    but does NOT normalise per-question dicts — Question.from_json needs the
    full original keys (``pairs``, ``answers``, ``question_type``, ...).
    """
    p = Path(path)
    with p.open("r", encoding="utf-8") as fh:
        data = json.load(fh)

    if isinstance(data, list):
        return {"questions": [q for q in data if isinstance(q, dict)]}

    if not isinstance(data, dict):
        return {"questions": []}

    if "sections" in data and isinstance(data["sections"], list):
        return data
    if "rounds" in data and isinstance(data["rounds"], list):
        return {"sections": [
            {"title": r.get("title") or "", "questions": r.get("questions") or []}
            for r in data["rounds"] if isinstance(r, dict)
        ]}
    if "questions" in data and isinstance(data["questions"], list):
        return {"questions": [q for q in data["questions"] if isinstance(q, dict)]}

    return {"questions": []}


# --- Public API ----------------------------------------------------------

def build_pptx(sections_or_questions, output_path, metadata=None, answers: str = "none") -> Path:
    if answers not in ANSWER_MODES:
        raise ValueError(f"answers must be one of {ANSWER_MODES}, got {answers!r}")

    p = Path(output_path)
    parent = p.parent
    if parent and not parent.exists():
        parent.mkdir(parents=True, exist_ok=True)
    elif parent and not parent.is_dir():
        raise OSError(f"Output parent exists and is not a directory: {parent}")

    meta = effective_metadata(metadata)
    raw_questions = _flatten_questions(sections_or_questions)
    questions: List[Question] = [Question.from_json(q) for q in raw_questions]
    total = len(questions)

    prs = Presentation()
    geom = _slide_geometry(prs)

    _add_cover_slide(prs, geom, meta)

    for i, q in enumerate(questions, start=1):
        _render_question_slide(prs, geom, meta, q, i, total, answers)
        if answers == "after":
            _add_answer_slide(prs, geom, meta, q)

    prs.save(str(p))
    print(f"PPTX deck created at: {p}")
    return p


# --- CLI -----------------------------------------------------------------

def main(argv: Optional[List[str]] = None) -> int:
    parser = argparse.ArgumentParser(description="Create a PPTX deck from a promptukit question bank.")
    parser.add_argument("-q", "--questions", required=True, help="Path to JSON question bank file")
    parser.add_argument("-o", "--output", required=True, help="Output .pptx filename")
    parser.add_argument("-m", "--metadata", "--setup", help="Path to JSON metadata/setup file", default=None)
    parser.add_argument(
        "--answers",
        choices=ANSWER_MODES,
        default="none",
        help="Answer reveal mode: none (default), after (extra slide per question), inline (answer at bottom of question slide)",
    )
    parser.add_argument("--save-questions", help="Write the normalized question artifact used for this deck", default=None)
    parser.add_argument("--save-setup", help="Write the effective metadata/setup artifact used for this deck", default=None)
    args = parser.parse_args(argv)

    questions_to_use = _read_question_bank(args.questions)

    meta = load_metadata_from_json(args.metadata) if args.metadata else None

    if args.save_questions:
        save_json_artifact(args.save_questions, questions_artifact(questions_to_use))
        print(f"Question artifact written to: {args.save_questions}")
    if args.save_setup:
        save_json_artifact(args.save_setup, effective_metadata(meta))
        print(f"Setup artifact written to: {args.save_setup}")

    try:
        build_pptx(questions_to_use, args.output, metadata=meta, answers=args.answers)
    except Exception as e:
        print(f"Error creating PPTX: {e}", file=sys.stderr)
        return 1
    return 0


if __name__ == "__main__":
    sys.exit(main())
